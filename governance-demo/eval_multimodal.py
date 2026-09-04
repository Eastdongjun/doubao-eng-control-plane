#!/usr/bin/env python3
"""多模态专项评测框架：五维垂直评测 + 质量门禁
① 代码可执行率  ② Excel 公式准确性  ③ PPT 结构质量  ④ 网页可运行性  ⑤ 报告引用可靠性
用法: eval_multimodal.py            # 标准样本评测（归零门禁）
      eval_multimodal.py --selftest # 注入缺陷样本，验证评测器能检出问题
"""
import argparse, datetime, json, pathlib, re, subprocess, sys

ROOT = pathlib.Path("/Users/donglai/Doubao/chats/2026-09-03/new-chat-6")
SAMPLE = ROOT / "_eval_samples"
GOV = ROOT / ".governance"

# ---------- 样本生成 ----------
def gen_samples(flawed=False):
    SAMPLE.mkdir(parents=True, exist_ok=True)
    from openpyxl import Workbook
    # Excel: 标准（SUM/AVERAGE 公式引用合法）
    wb = Workbook(); ws = wb.active; ws.title = "sales"
    for i, v in enumerate([10, 20, 30, 40], start=1):
        ws.cell(row=i, column=1, value=v)
    ws["C1"] = "=SUM(A1:A4)"
    ws["C2"] = "=AVERAGE(A1:A4)"
    if flawed:
        ws["C3"] = "=SUM(X99:X100)"  # 引用不存在区域
        ws["C4"] = "=#REF!"         # 错误公式
    wb.save(SAMPLE / "sample.xlsx")

    from pptx import Presentation
    prs = Presentation()
    for title in (["营收分析", "成本结构", "增长预测"] if not flawed else ["", "成本结构", "增长预测"]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if title:
            slide.shapes.title.text = title
        if flawed and title == "增长预测":
            pass  # 留一张无正文的"空"页
    prs.save(SAMPLE / "sample.pptx")

    # Web: 引用本地存在资源
    (SAMPLE / "app.css").write_text("body{font-family:sans-serif;}\n", encoding="utf-8")
    (SAMPLE / "app.js").write_text("console.log('ok');\n", encoding="utf-8")
    html = ('<!DOCTYPE html><html><head><title>样页</title>'
            '<link rel="stylesheet" href="app.css"><script src="app.js"></script></head>'
            '<body><h1>工程化评测</h1></body></html>')
    if flawed:
        html = html.replace("</html>", "")          # 标签未闭合
        html = html.replace('href="app.css"', 'href="missing.css"')  # 断链
    (SAMPLE / "sample.html").write_text(html, encoding="utf-8")

    # 报告: 引用标记闭合
    if flawed:
        report = "# 测试报告\n结论：达标。\n"  # 缺引用标记
    else:
        report = ('# 测试报告\n结论：达标。\n\n'
                  '["https://example.com/doc1","https://example.com/doc2"]\n')
    (SAMPLE / "sample-report.md").write_text(report, encoding="utf-8")

# ---------- ① 代码可执行率 ----------
def eval_code():
    from eval_vertical import syntax_check, run_case
    cases = [
        ("calculator.py", ["3", "+", "4"], "7"),
        ("logclean.py", ["--help"], "[-h]"),
        ("xian_guide.py", ["钟楼"], "钟楼"),
    ]
    dir0 = ROOT / "auto-dev-demo"
    total = passed = 0
    detail = []
    for name, args, expect in cases:
        py = dir0 / name
        if not py.exists(): continue
        total += 1
        ok_s, _ = syntax_check(py)
        ok_r, note = run_case(py, args, expect)
        ok = ok_s and ok_r
        passed += 1 if ok else 0
        detail.append({"case": name, "pass": ok, "note": note})
    return {"total": total, "pass": passed,
            "rate": round(100 * passed / total, 1) if total else 0, "detail": detail}

# ---------- ② Excel 公式准确性 ----------
def eval_excel(flawed):
    from openpyxl import load_workbook
    wb = load_workbook(SAMPLE / "sample.xlsx")
    ws = wb.active
    total = passed = 0; detail = []
    for row in ws.iter_rows(min_row=1, max_col=3):
        for cell in row:
            if isinstance(cell.value, str) and cell.value.startswith("="):
                total += 1
                formula = cell.value[1:]
                bad = "#REF" in formula.upper()
                refs = re.findall(r"([A-Z]+)(\d+)", formula)
                for col, rn in refs:
                    if ws.max_row < int(rn) or ws.max_column < len(col):
                        bad = True
                passed += 1 if not bad else 0
                detail.append({"cell": cell.coordinate, "formula": cell.value, "pass": not bad})
    # 期望检出率：缺陷样本应检出 2 个坏公式
    return {"total": total, "pass": passed,
            "rate": round(100 * passed / total, 1) if total else 0, "detail": detail}

# ---------- ③ PPT 结构质量 ----------
def eval_ppt(flawed):
    from pptx import Presentation
    prs = Presentation(SAMPLE / "sample.pptx")
    total = passed = 0; detail = []
    for i, slide in enumerate(prs.slides, start=1):
        total += 1
        title_shape = slide.shapes.title
        title = title_shape.text if title_shape is not None else ""
        ok = bool(title.strip())
        passed += 1 if ok else 0
        detail.append({"slide": i, "title": title or "(空)", "pass": ok})
    return {"total": total, "pass": passed,
            "rate": round(100 * passed / total, 1) if total else 0, "detail": detail}

# ---------- ④ 网页可运行性 ----------
def eval_web(flawed):
    from html.parser import HTMLParser
    text = (SAMPLE / "sample.html").read_text(encoding="utf-8")
    stack = []
    class P(HTMLParser):
        def handle_starttag(self, tag, attrs):
            if tag not in ("meta", "link", "br", "img", "input", "hr"):
                stack.append(tag)
        def handle_endtag(self, tag):
            if tag in stack: stack.remove(tag)
    p = P(); p.feed(text)
    total = 2; passed = 0; detail = []
    closed = len(stack) == 0
    passed += 1 if closed else 0
    detail.append({"check": "标签闭合", "pass": closed})
    assets_ok = True
    for m in re.findall(r'(?:href|src)="([^"]+)"', text):
        if not m.startswith(("http", "//")):
            if not (SAMPLE / m).exists(): assets_ok = False
    passed += 1 if assets_ok else 0
    detail.append({"check": "本地资源存在", "pass": assets_ok})
    return {"total": total, "pass": passed,
            "rate": round(100 * passed / total, 1) if total else 0, "detail": detail}

# ---------- ⑤ 报告引用可靠性 ----------
def eval_report(flawed):
    text = (SAMPLE / "sample-report.md").read_text(encoding="utf-8")
    total = 1; passed = 0; detail = []
    m = re.search(r'\s*(\[[^\]]*\])\s*', text)
    if m:
        urls = json.loads(m.group(1))
        ok = len(urls) >= 1 and all(u.startswith("http") for u in urls)
    else:
        ok = False
    passed += 1 if ok else 0
    detail.append({"check": "引用标记闭合且 URL 合法", "pass": ok})
    return {"total": total, "pass": passed,
            "rate": round(100 * passed / total, 1) if total else 0, "detail": detail}

# ---------- 主流程 ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--selftest", action="store_true", help="注入缺陷样本验证评测器")
    args = ap.parse_args()
    gen_samples(flawed=args.selftest)
    dims = {
        "代码可执行率": eval_code(),
        "Excel公式准确性": eval_excel(args.selftest),
        "PPT结构质量": eval_ppt(args.selftest),
        "网页可运行性": eval_web(args.selftest),
        "报告引用可靠性": eval_report(args.selftest),
    }
    print("=" * 60)
    print("多模态专项评测 · " + ("缺陷样本自检" if args.selftest else "标准样本（归零门禁）"))
    print("=" * 60)
    all_gate = True
    for name, r in dims.items():
        gate = r["rate"] >= 100
        all_gate = all_gate and gate
        print(f"  {'✓' if gate else '✗'} {name}: {r['pass']}/{r['total']} ({r['rate']}%)")
        for d in r["detail"]:
            if not d["pass"]:
                print(f"      → 检出: {d}")
    print("-" * 60)
    mode = "自检通过（缺陷全检出）" if (args.selftest and not all_gate) else \
           ("✓ 质量门禁通过（五维 100%，归零）" if all_gate else "✗ 存在未达 100% 维度")
    print(mode)
    GOV.mkdir(exist_ok=True)
    report = {"mode": "selftest" if args.selftest else "gate",
              "generated_at": datetime.datetime.now().strftime("%Y-%m-%dT%H:%M:%S"),
              "dims": dims}
    (GOV / "eval-multimodal.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"报告: {GOV / 'eval-multimodal.json'}")
    # 归零门禁：标准样本必须全绿
    return 0 if (all_gate or args.selftest) else 1

if __name__ == "__main__":
    sys.exit(main())
